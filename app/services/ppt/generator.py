"""
PowerPoint presentation generation service.

Responsible for creating actual PPTX files from validated presentation plans.
"""

import io
import logging
import os
import tempfile
from pathlib import Path
from typing import Optional, List, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from app.schemas.presentations import PresentationPlan, Slide, SlideLayout, Theme, Chart, Table

logger = logging.getLogger("app.ppt.generator")


class PPTGenerator:
    """Core PowerPoint file generator using python-pptx."""

    # Theme color palettes
    THEMES = {
        Theme.modern: {
            "bg_color": RGBColor(255, 255, 255),
            "title_color": RGBColor(0, 51, 102),
            "accent_color": RGBColor(0, 153, 255),
            "text_color": RGBColor(51, 51, 51),
        },
        Theme.professional: {
            "bg_color": RGBColor(245, 245, 245),
            "title_color": RGBColor(44, 62, 80),
            "accent_color": RGBColor(41, 128, 185),
            "text_color": RGBColor(52, 73, 94),
        },
        Theme.minimal: {
            "bg_color": RGBColor(255, 255, 255),
            "title_color": RGBColor(0, 0, 0),
            "accent_color": RGBColor(128, 128, 128),
            "text_color": RGBColor(80, 80, 80),
        },
        Theme.education: {
            "bg_color": RGBColor(240, 248, 255),
            "title_color": RGBColor(25, 25, 112),
            "accent_color": RGBColor(255, 140, 0),
            "text_color": RGBColor(64, 64, 64),
        },
        Theme.dark: {
            "bg_color": RGBColor(30, 30, 30),
            "title_color": RGBColor(255, 255, 255),
            "accent_color": RGBColor(100, 200, 255),
            "text_color": RGBColor(200, 200, 200),
        },
    }

    # Default slide dimensions (standard 16:9)
    SLIDE_WIDTH = Inches(10)
    SLIDE_HEIGHT = Inches(7.5)

    # Margins
    LEFT_MARGIN = Inches(0.5)
    RIGHT_MARGIN = Inches(0.5)
    TOP_MARGIN = Inches(0.5)
    BOTTOM_MARGIN = Inches(0.5)

    def __init__(self, presentation_plan: PresentationPlan):
        """Initialize generator with presentation plan."""
        self.plan = presentation_plan
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT
        self.theme_colors = self.THEMES.get(presentation_plan.theme, self.THEMES[Theme.modern])

    def generate(self) -> bytes:
        """Generate PPTX and return as bytes."""
        try:
            logger.info(
                "Generating PPT: '%s' with %d slides, theme=%s",
                self.plan.title,
                len(self.plan.slides),
                self.plan.theme,
            )

            for slide_spec in self.plan.slides:
                self._add_slide(slide_spec)

            # Return as bytes
            output = io.BytesIO()
            self.prs.save(output)
            output.seek(0)
            pptx_bytes = output.getvalue()

            logger.info("PPT generated successfully: %d bytes", len(pptx_bytes))
            return pptx_bytes

        except Exception as exc:
            logger.exception("PPT generation failed: %s", exc)
            raise

    def _add_slide(self, slide_spec: Slide) -> None:
        """Add a single slide to the presentation."""
        layout_type = slide_spec.layout

        if layout_type == SlideLayout.title:
            self._add_title_slide(slide_spec)
        elif layout_type == SlideLayout.section:
            self._add_section_slide(slide_spec)
        elif layout_type == SlideLayout.content:
            self._add_content_slide(slide_spec)
        elif layout_type == SlideLayout.title_and_content:
            self._add_title_and_content_slide(slide_spec)
        elif layout_type == SlideLayout.two_column:
            self._add_two_column_slide(slide_spec)
        elif layout_type == SlideLayout.quote:
            self._add_quote_slide(slide_spec)
        elif layout_type == SlideLayout.code:
            self._add_code_slide(slide_spec)
        elif layout_type == SlideLayout.table:
            self._add_table_slide(slide_spec)
        elif layout_type == SlideLayout.chart:
            self._add_chart_slide(slide_spec)
        elif layout_type == SlideLayout.conclusion:
            self._add_conclusion_slide(slide_spec)
        else:
            logger.warning("Unknown layout type: %s, using content", layout_type)
            self._add_content_slide(slide_spec)

    def _add_title_slide(self, slide_spec: Slide) -> None:
        """Add title slide (cover page)."""
        blank_slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme_colors["bg_color"]

        # Title
        title_box = slide.shapes.add_textbox(
            self.LEFT_MARGIN,
            Inches(2.5),
            self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
            Inches(1.5),
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = title_frame.paragraphs[0]
        p.text = slide_spec.title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["title_color"]
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if slide_spec.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                self.LEFT_MARGIN,
                Inches(4.2),
                self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
                Inches(1),
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            p = subtitle_frame.paragraphs[0]
            p.text = slide_spec.subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = self.theme_colors["accent_color"]
            p.alignment = PP_ALIGN.CENTER

        # Author
        author_box = slide.shapes.add_textbox(
            self.LEFT_MARGIN,
            Inches(6.5),
            self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
            Inches(0.7),
        )
        author_frame = author_box.text_frame
        p = author_frame.paragraphs[0]
        p.text = self.plan.author or "VedaApex"
        p.font.size = Pt(16)
        p.font.color.rgb = self.theme_colors["text_color"]
        p.alignment = PP_ALIGN.CENTER

    def _add_section_slide(self, slide_spec: Slide) -> None:
        """Add section divider slide."""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme_colors["accent_color"]

        # Title
        title_box = slide.shapes.add_textbox(
            self.LEFT_MARGIN,
            Inches(3),
            self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
            Inches(1.5),
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = title_frame.paragraphs[0]
        p.text = slide_spec.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, slide_spec: Slide) -> None:
        """Add standard content slide with title and bullet points."""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme_colors["bg_color"]

        # Title
        title_box = slide.shapes.add_textbox(
            self.LEFT_MARGIN,
            self.TOP_MARGIN,
            self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
            Inches(0.8),
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = slide_spec.title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["title_color"]

        # Content area
        content_top = self.TOP_MARGIN + Inches(1)
        content_height = self.SLIDE_HEIGHT - content_top - self.BOTTOM_MARGIN

        # Bullets
        if slide_spec.bullets:
            bullets_box = slide.shapes.add_textbox(
                self.LEFT_MARGIN + Inches(0.3),
                content_top,
                self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN - Inches(0.3),
                content_height * 0.6,
            )
            text_frame = bullets_box.text_frame
            text_frame.word_wrap = True

            for i, bullet in enumerate(slide_spec.bullets):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = self._truncate_text(bullet, 200)
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = self.theme_colors["text_color"]
                p.space_before = Pt(6)
                p.space_after = Pt(6)

        # Paragraphs
        if slide_spec.paragraphs:
            para_top = content_top + Inches(2.5)
            para_box = slide.shapes.add_textbox(
                self.LEFT_MARGIN,
                para_top,
                self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
                content_height * 0.4,
            )
            text_frame = para_box.text_frame
            text_frame.word_wrap = True

            for i, para in enumerate(slide_spec.paragraphs):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = self._truncate_text(para, 500)
                p.font.size = Pt(14)
                p.font.color.rgb = self.theme_colors["text_color"]
                p.space_after = Pt(12)

    def _add_title_and_content_slide(self, slide_spec: Slide) -> None:
        """Add slide with title and mixed content."""
        self._add_content_slide(slide_spec)

    def _add_two_column_slide(self, slide_spec: Slide) -> None:
        """Add two-column layout slide."""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme_colors["bg_color"]

        # Title
        title_box = slide.shapes.add_textbox(
            self.LEFT_MARGIN,
            self.TOP_MARGIN,
            self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
            Inches(0.8),
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = slide_spec.title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["title_color"]

        # Left column
        if slide_spec.bullets:
            left_box = slide.shapes.add_textbox(
                self.LEFT_MARGIN,
                self.TOP_MARGIN + Inches(1),
                (self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN) / 2 - Inches(0.2),
                Inches(5.5),
            )
            text_frame = left_box.text_frame
            text_frame.word_wrap = True

            for i, bullet in enumerate(slide_spec.bullets[: len(slide_spec.bullets) // 2 + 1]):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = self._truncate_text(bullet, 150)
                p.level = 0
                p.font.size = Pt(16)
                p.font.color.rgb = self.theme_colors["text_color"]

        # Right column
        if len(slide_spec.bullets) > len(slide_spec.bullets) // 2 + 1:
            right_box = slide.shapes.add_textbox(
                self.LEFT_MARGIN + (self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN) / 2 + Inches(0.2),
                self.TOP_MARGIN + Inches(1),
                (self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN) / 2 - Inches(0.2),
                Inches(5.5),
            )
            text_frame = right_box.text_frame
            text_frame.word_wrap = True

            for i, bullet in enumerate(slide_spec.bullets[len(slide_spec.bullets) // 2 + 1 :]):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = self._truncate_text(bullet, 150)
                p.level = 0
                p.font.size = Pt(16)
                p.font.color.rgb = self.theme_colors["text_color"]

    def _add_quote_slide(self, slide_spec: Slide) -> None:
        """Add quote/inspiration slide."""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme_colors["bg_color"]

        # Quote text
        if slide_spec.quote:
            quote_box = slide.shapes.add_textbox(
                self.LEFT_MARGIN + Inches(0.5),
                Inches(2.5),
                self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN - Inches(1),
                Inches(2.5),
            )
            text_frame = quote_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = f'"{self._truncate_text(slide_spec.quote, 300)}"'
            p.font.size = Pt(32)
            p.font.italic = True
            p.font.color.rgb = self.theme_colors["accent_color"]
            p.alignment = PP_ALIGN.CENTER

        # Author
        if slide_spec.author:
            author_box = slide.shapes.add_textbox(
                self.LEFT_MARGIN + Inches(0.5),
                Inches(5.2),
                self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN - Inches(1),
                Inches(1),
            )
            text_frame = author_box.text_frame
            p = text_frame.paragraphs[0]
            p.text = f"— {slide_spec.author}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.theme_colors["text_color"]
            p.alignment = PP_ALIGN.RIGHT

    def _add_code_slide(self, slide_spec: Slide) -> None:
        """Add code snippet slide."""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme_colors["bg_color"]

        # Title
        title_box = slide.shapes.add_textbox(
            self.LEFT_MARGIN,
            self.TOP_MARGIN,
            self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
            Inches(0.7),
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = slide_spec.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["title_color"]

        # Code box
        if slide_spec.code:
            code_box = slide.shapes.add_textbox(
                self.LEFT_MARGIN,
                self.TOP_MARGIN + Inches(1),
                self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
                Inches(5.5),
            )
            text_frame = code_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = self._truncate_text(slide_spec.code, 1500)
            p.font.name = "Courier New"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0, 100, 0)

    def _add_table_slide(self, slide_spec: Slide) -> None:
        """Add table slide."""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme_colors["bg_color"]

        # Title
        title_box = slide.shapes.add_textbox(
            self.LEFT_MARGIN,
            self.TOP_MARGIN,
            self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
            Inches(0.7),
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = slide_spec.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["title_color"]

        # Table
        if slide_spec.table:
            table_spec = slide_spec.table
            rows = len(table_spec.rows) + 1
            cols = len(table_spec.headers)

            table_width = self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN
            table_height = Inches(4.5)

            left = self.LEFT_MARGIN
            top = self.TOP_MARGIN + Inches(1.2)

            table_shape = slide.shapes.add_table(rows, cols, left, top, table_width, table_height)
            table_obj = table_shape.table

            # Headers
            for col_idx, header in enumerate(table_spec.headers):
                cell = table_obj.cell(0, col_idx)
                cell.text = header[:50]  # Truncate
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.theme_colors["accent_color"]
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)

            # Rows
            for row_idx, row_spec in enumerate(table_spec.rows, start=1):
                for col_idx, cell_text in enumerate(row_spec.cells):
                    cell = table_obj.cell(row_idx, col_idx)
                    cell.text = cell_text[:100]  # Truncate
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(11)

    def _add_chart_slide(self, slide_spec: Slide) -> None:
        """Add chart slide (basic implementation)."""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme_colors["bg_color"]

        # Title
        title_box = slide.shapes.add_textbox(
            self.LEFT_MARGIN,
            self.TOP_MARGIN,
            self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
            Inches(0.7),
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = slide_spec.title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.theme_colors["title_color"]

        # Placeholder message
        if slide_spec.chart:
            chart_spec = slide_spec.chart
            message_box = slide.shapes.add_textbox(
                self.LEFT_MARGIN + Inches(1),
                Inches(3),
                self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN - Inches(2),
                Inches(3),
            )
            text_frame = message_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = (
                f"Chart: {chart_spec.type.value.upper()}\n"
                f"Title: {chart_spec.title}\n"
                f"Categories: {', '.join(chart_spec.categories[:5])}"
            )
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme_colors["text_color"]

    def _add_conclusion_slide(self, slide_spec: Slide) -> None:
        """Add conclusion/closing slide."""
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme_colors["accent_color"]

        # Title
        title_box = slide.shapes.add_textbox(
            self.LEFT_MARGIN,
            Inches(2),
            self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
            Inches(2),
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = title_frame.paragraphs[0]
        p.text = slide_spec.title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if slide_spec.subtitle:
            subtitle_box = slide.shapes.add_textbox(
                self.LEFT_MARGIN,
                Inches(4.5),
                self.SLIDE_WIDTH - self.LEFT_MARGIN - self.RIGHT_MARGIN,
                Inches(1.5),
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            p = subtitle_frame.paragraphs[0]
            p.text = slide_spec.subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

    @staticmethod
    def _truncate_text(text: str, max_length: int) -> str:
        """Truncate text to prevent overflow, with ellipsis."""
        if len(text) <= max_length:
            return text
        return text[:max_length - 3] + "..."
