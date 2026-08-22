from app.utils.time import utcnow

import os
import json
import uuid
import httpx
import logging
from pathlib import Path
from typing import Any
from sqlmodel import Session
from fastapi import HTTPException
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt
from docx import Document
from docx.shared import Inches as DocxInches
from openpyxl import Workbook
from openpyxl.drawing.image import Image as ExcelImage
from app.core.config import settings
from app.services.attachments.processor import AttachmentProcessor
from app.storage.storage_manager import storage_manager

logger = logging.getLogger(__name__)


class DocumentService:
    """Service to programmatically generate PPT, Word, and Excel files using multiple LLM providers."""

    TEXT_PROVIDER_ALIASES = {
        "auto",
        "text",
        "text_generation",
        "text-generation",
        "document",
        "document_compiler",
        "compiler",
    }

    @staticmethod
    def _clean_json_response(text: str) -> str:
        """Strip markdown markers if returned by the LLM."""
        text = text.strip()
        if text.startswith("```"):
            lines = text.split("\n")
            if lines[0].startswith("```"):
                lines = lines[1:]
            if lines[-1].startswith("```"):
                lines = lines[:-1]
            text = "\n".join(lines).strip()
        return text

    @staticmethod
    def _parse_json_safe(cleaned: str) -> dict:
        """Parse JSON from cleaned LLM output with fallback extraction."""
        try:
            return json.loads(cleaned)
        except Exception as e:
            # Fallback parsing helper if JSON starts/ends with extra text
            try:
                start = cleaned.find("{")
                end = cleaned.rfind("}")
                if start != -1 and end != -1:
                    return json.loads(cleaned[start : end + 1])
                start_arr = cleaned.find("[")
                end_arr = cleaned.rfind("]")
                if start_arr != -1 and end_arr != -1:
                    return json.loads(cleaned[start_arr : end_arr + 1])
            except Exception:
                pass
            raise HTTPException(
                status_code=500,
                detail=f"Failed to parse structured JSON from compiler: {e}",
            )

    @staticmethod
    def _normalize_provider(provider: str | None) -> str:
        if not provider:
            return "auto"
        normalized = provider.strip().lower()
        return normalized or "auto"

    @staticmethod
    def _compiler_key() -> str:
        return (
            settings.DOCUMENT_COMPILER_KEY
            or settings.DOCUMENT_GENERATION_API_KEY
            or settings.TEXT_GENERATION_API_KEY
            or settings.OPENAI_API_KEY
            or settings.GEMINI_API_KEY
            or settings.GEMINI_API_KEY_TIER1
            or ""
        ).strip()

    @staticmethod
    def _save_generated_file(filename: str, file_path: str, base_url: str) -> str:
        """Store a generated file through the app storage manager and return a full public URL."""
        if not os.path.exists(file_path):
            raise HTTPException(status_code=500, detail=f"Generated file not found: {file_path}")
        public_path = storage_manager.upload_file(file_path, filename)
        if public_path.startswith("http"):
            return public_path
        return f"{base_url.rstrip('/')}{public_path}"

    @staticmethod
    def _clean_text(value: Any, fallback: str = "") -> str:
        if value is None:
            return fallback
        text = str(value).strip()
        return text or fallback

    @staticmethod
    def _pdf_safe_text(value: Any) -> str:
        return str(value or "").encode("latin-1", errors="replace").decode("latin-1")

    @staticmethod
    def _attachment_rows(attachments: list[dict[str, Any]] | None) -> list[dict[str, str]]:
        rows: list[dict[str, str]] = []
        for item in attachments or []:
            filename = DocumentService._clean_text(item.get("filename"), "attachment")
            mime_type = DocumentService._clean_text(item.get("mime_type"), "application/octet-stream")
            extension = Path(filename).suffix.lower()
            size = item.get("size") or 0
            kind = "file"
            if mime_type.startswith("image/") or extension in {".jpg", ".jpeg", ".png", ".webp", ".gif"}:
                kind = "image"
            elif mime_type.startswith("video/") or extension in {".mp4", ".mov", ".webm", ".avi", ".mkv"}:
                kind = "video"
            elif (
                mime_type.startswith("text/")
                or "document" in mime_type
                or "pdf" in mime_type
                or "spreadsheet" in mime_type
                or "presentation" in mime_type
                or extension in {".csv", ".docx", ".json", ".pdf", ".pptx", ".txt", ".xlsx"}
            ):
                kind = "document"
            rows.append(
                {
                    "filename": filename,
                    "mime_type": mime_type,
                    "kind": kind,
                    "size": str(size),
                    "path": str(item.get("temp_path") or item.get("path") or ""),
                    "url": str(item.get("url") or item.get("source_url") or ""),
                }
            )
        return rows

    @staticmethod
    def _build_attachment_context(attachments: list[dict[str, Any]] | None) -> str:
        parts: list[str] = []
        for index, row in enumerate(DocumentService._attachment_rows(attachments), start=1):
            parts.append(
                f"{index}. {row['filename']} ({row['kind']}, {row['mime_type']}, {row['size']} bytes)"
            )
            path = row["path"]
            if path and Path(path).exists():
                text = AttachmentProcessor.extract_document_text(path, row["mime_type"])
                if text:
                    parts.append(f"Extracted content:\n{text[:6000]}")
        return "\n\n".join(parts).strip()

    @staticmethod
    def prepare_prompt_with_attachments(
        prompt: str,
        attachments: list[dict[str, Any]] | None = None,
        attachment_urls: list[str] | None = None,
    ) -> str:
        prompt_text = prompt.strip()
        context = DocumentService._build_attachment_context(attachments)
        url_lines = [
            f"- {url.strip()}"
            for url in attachment_urls or []
            if isinstance(url, str) and url.strip()
        ]

        additions: list[str] = []
        if context:
            additions.append("Uploaded attachment context:\n" + context)
        if url_lines:
            additions.append("Referenced attachment URLs:\n" + "\n".join(url_lines))

        if not additions:
            return prompt_text

        return (
            f"{prompt_text}\n\n"
            "Use the following attachment details as source material where relevant. "
            "If an uploaded video is listed, mention it as an attachment/reference unless its transcript is present.\n\n"
            + "\n\n".join(additions)
        )

    @staticmethod
    def _fallback_title(prompt: str, default: str) -> str:
        words = " ".join((prompt or "").strip().split())
        if not words:
            return default
        return words[:80].rstrip(" .,;:-") or default

    @staticmethod
    def _fallback_structured_data(prompt: str, kind: str) -> dict:
        title = DocumentService._fallback_title(prompt, "Generated Document")
        summary = prompt.strip() or "Generated from the supplied prompt and attachments."

        if kind == "ppt":
            return {
                "slides": [
                    {"title": title, "bullet_points": ["Overview", summary[:220]]},
                    {"title": "Key Points", "bullet_points": [summary[:220], "Details are based on the provided prompt and attachments."]},
                    {"title": "Next Steps", "bullet_points": ["Review the generated content.", "Update or expand sections as needed."]},
                ]
            }

        if kind == "excel":
            return {
                "sheet_title": title[:30] or "Data",
                "headers": ["Item", "Details", "Source"],
                "rows": [
                    ["Prompt", summary[:320], "User request"],
                    ["Generated At", utcnow().strftime("%Y-%m-%d %H:%M UTC"), "System"],
                ],
            }

        sections = [
            {
                "heading": "Overview",
                "paragraphs": [summary[:900]],
            },
            {
                "heading": "Details",
                "paragraphs": ["This file was generated from the supplied prompt and available attachment context."],
            },
        ]
        data = {"title": title, "sections": sections}
        if kind == "pdf":
            data["subtitle"] = "Generated document"
        return data

    @staticmethod
    def _normalize_slides(data: dict, prompt: str) -> list[dict[str, Any]]:
        slides = data.get("slides")
        if slides is None and isinstance(data.get("sections"), list):
            slides = [
                {
                    "title": section.get("heading", "Section"),
                    "bullet_points": section.get("paragraphs", []),
                }
                for section in data["sections"]
                if isinstance(section, dict)
            ]
        if not isinstance(slides, list) or not slides:
            slides = DocumentService._fallback_structured_data(prompt, "ppt")["slides"]

        normalized = []
        for slide_data in slides:
            if not isinstance(slide_data, dict):
                slide_data = {"title": "Slide", "bullet_points": [str(slide_data)]}
            points = slide_data.get("bullet_points") or slide_data.get("bullets") or slide_data.get("points") or []
            if isinstance(points, str):
                points = [points]
            normalized.append(
                {
                    "title": DocumentService._clean_text(slide_data.get("title"), "Untitled Slide"),
                    "bullet_points": [DocumentService._clean_text(point) for point in points if DocumentService._clean_text(point)],
                }
            )
        return normalized

    @staticmethod
    def _normalize_sections(data: dict, prompt: str, kind: str = "word") -> dict:
        fallback = DocumentService._fallback_structured_data(prompt, kind)
        title = DocumentService._clean_text(data.get("title"), fallback["title"])
        subtitle = DocumentService._clean_text(data.get("subtitle"), fallback.get("subtitle", ""))
        sections = data.get("sections")

        if not isinstance(sections, list) or not sections:
            sections = fallback["sections"]

        normalized_sections = []
        for section in sections:
            if not isinstance(section, dict):
                section = {"heading": "Section", "paragraphs": [str(section)]}
            paragraphs = section.get("paragraphs") or section.get("content") or []
            if isinstance(paragraphs, str):
                paragraphs = [paragraphs]
            normalized_sections.append(
                {
                    "heading": DocumentService._clean_text(section.get("heading"), "Section"),
                    "paragraphs": [DocumentService._clean_text(item) for item in paragraphs if DocumentService._clean_text(item)],
                }
            )

        result = {"title": title, "sections": normalized_sections}
        if kind == "pdf":
            result["subtitle"] = subtitle
        return result

    @staticmethod
    def _normalize_sheet(data: dict, prompt: str) -> dict:
        fallback = DocumentService._fallback_structured_data(prompt, "excel")
        sheet_title = DocumentService._clean_text(data.get("sheet_title") or data.get("title"), fallback["sheet_title"])[:30]
        headers = data.get("headers") or fallback["headers"]
        rows = data.get("rows") or fallback["rows"]
        if not isinstance(headers, list) or not headers:
            headers = fallback["headers"]
        if not isinstance(rows, list) or not rows:
            rows = fallback["rows"]
        return {
            "sheet_title": sheet_title or "Data",
            "headers": [DocumentService._clean_text(header, "Column") for header in headers],
            "rows": [
                [DocumentService._clean_text(cell) for cell in row] if isinstance(row, list) else [DocumentService._clean_text(row)]
                for row in rows
            ],
        }

    @staticmethod
    def _image_attachments(attachments: list[dict[str, Any]] | None) -> list[dict[str, str]]:
        images = []
        for row in DocumentService._attachment_rows(attachments):
            if row["kind"] != "image":
                continue
            path = row["path"]
            if path and Path(path).exists():
                images.append(row)
        return images

    @staticmethod
    def _non_image_attachment_rows(attachments: list[dict[str, Any]] | None) -> list[dict[str, str]]:
        return [row for row in DocumentService._attachment_rows(attachments) if row["kind"] != "image"]

    # =========================================================================
    #  Provider-specific content generators
    # =========================================================================

    @staticmethod
    async def _generate_via_groq(prompt: str, system_prompt: str, tier: int = 1) -> dict:
        """Generate structured JSON content using Groq (llama-3.3-70b-versatile)."""
        from app.services.providers.groq_provider import GroqProvider

        result = await GroqProvider.run_model(
            "llama-3.3-70b-versatile",
            {
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt},
                ],
                "max_tokens": 4096,
                "temperature": 0.7,
            },
            tier,
        )

        # Extract content from OpenAI-compatible response
        if isinstance(result, dict) and "choices" in result:
            content = result["choices"][0]["message"]["content"]
        elif isinstance(result, str):
            content = result
        else:
            raise HTTPException(
                status_code=500, detail=f"Unexpected Groq response format: {result}"
            )

        cleaned = DocumentService._clean_json_response(content)
        return DocumentService._parse_json_safe(cleaned)

    @staticmethod
    async def _generate_via_ollama(prompt: str, system_prompt: str, tier: int = 1) -> dict:
        """Generate structured JSON content using Ollama (llama3)."""
        from app.services.providers.ollama_provider import OllamaProvider

        result = await OllamaProvider.run_model(
            "llama3",
            {
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt},
                ],
                "temperature": 0.7,
            },
            tier,
        )

        # Extract content from OpenAI-compatible response
        if isinstance(result, dict) and "choices" in result:
            content = result["choices"][0]["message"]["content"]
        elif isinstance(result, str):
            content = result
        else:
            raise HTTPException(
                status_code=500, detail=f"Unexpected Ollama response format: {result}"
            )

        cleaned = DocumentService._clean_json_response(content)
        return DocumentService._parse_json_safe(cleaned)

    @staticmethod
    async def _generate_via_compiler(prompt: str, system_prompt: str) -> dict:
        """Generate structured JSON content using the dedicated DOCUMENT_COMPILER_KEY (Gemini or OpenAI)."""
        api_key = DocumentService._compiler_key()
        if not api_key:
            raise HTTPException(
                status_code=501,
                detail="Document compilation feature is currently inactive. Configure DOCUMENT_COMPILER_KEY, DOCUMENT_GENERATION_API_KEY, or TEXT_GENERATION_API_KEY.",
            )

        # 1. If key starts with sk-, treat as OpenAI API key
        if api_key.startswith("sk-"):
            url = "https://api.openai.com/v1/chat/completions"
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            }
            payload = {
                "model": "gpt-4o-mini",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt},
                ],
                "response_format": {"type": "json_object"},
            }
            try:
                async with httpx.AsyncClient(timeout=60.0) as client:
                    response = await client.post(url, headers=headers, json=payload)
                    if response.status_code != 200:
                        raise Exception(f"OpenAI error response: {response.text}")
                    res_data = response.json()
                    content = res_data["choices"][0]["message"]["content"]
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"OpenAI compilation failed: {e}")

        # 2. Otherwise, treat as Google Gemini API Key
        else:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}"
            payload = {
                "contents": [
                    {
                        "parts": [
                            {"text": f"System Guidelines: {system_prompt}"},
                            {"text": f"User Prompt: {prompt}"},
                        ]
                    }
                ],
                "generationConfig": {"responseMimeType": "application/json"},
            }
            try:
                async with httpx.AsyncClient(timeout=60.0) as client:
                    response = await client.post(url, json=payload)
                    if response.status_code != 200:
                        raise Exception(f"Gemini error response: {response.text}")
                    res_data = response.json()
                    content = res_data["candidates"][0]["content"]["parts"][0]["text"]
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"Gemini compilation failed: {e}")

        cleaned = DocumentService._clean_json_response(content)
        return DocumentService._parse_json_safe(cleaned)

    @staticmethod
    async def _generate_via_text_service(
        prompt: str,
        system_prompt: str,
        provider: str = "auto",
        tier: int = 1,
    ) -> dict:
        """Generate structured JSON using the normal text generation provider chain."""
        from app.services.ai_service import AIToolsService

        provider_name = DocumentService._normalize_provider(provider)
        if provider_name in DocumentService.TEXT_PROVIDER_ALIASES:
            provider_name = "auto"

        result = await AIToolsService.generate_text(
            prompt=prompt,
            system_prompt=system_prompt,
            tier=tier,
            provider=provider_name,
        )

        if isinstance(result, dict) and "choices" in result:
            content = result["choices"][0]["message"]["content"]
        elif isinstance(result, str):
            content = result
        else:
            content = json.dumps(result)

        cleaned = DocumentService._clean_json_response(content)
        return DocumentService._parse_json_safe(cleaned)

    @staticmethod
    async def _generate_content(
        prompt: str,
        system_prompt: str,
        provider: str = "auto",
        tier: int = 1,
        kind: str = "word",
    ) -> dict:
        """Route to the correct LLM backend based on provider."""
        provider_name = DocumentService._normalize_provider(provider)
        errors: list[str] = []

        if provider_name in {"document_compiler", "compiler"}:
            if DocumentService._compiler_key():
                try:
                    return await DocumentService._generate_via_compiler(prompt, system_prompt)
                except Exception as exc:
                    errors.append(f"document_compiler: {exc}")
            provider_name = "auto"

        if provider_name == "auto":
            if DocumentService._compiler_key():
                try:
                    return await DocumentService._generate_via_compiler(prompt, system_prompt)
                except Exception as exc:
                    errors.append(f"document_compiler: {exc}")
            try:
                return await DocumentService._generate_via_text_service(
                    prompt, system_prompt, "auto", tier
                )
            except Exception as exc:
                errors.append(f"text:auto: {exc}")
        elif provider_name == "groq":
            try:
                return await DocumentService._generate_via_groq(prompt, system_prompt, tier)
            except Exception as exc:
                errors.append(f"groq: {exc}")
        elif provider_name == "ollama":
            try:
                return await DocumentService._generate_via_ollama(prompt, system_prompt, tier)
            except Exception as exc:
                errors.append(f"ollama: {exc}")
        else:
            try:
                return await DocumentService._generate_via_text_service(
                    prompt, system_prompt, provider_name, tier
                )
            except Exception as exc:
                errors.append(f"{provider_name}: {exc}")

        logger.warning("Document content generation fell back to local structure: %s", "; ".join(errors))
        return DocumentService._fallback_structured_data(prompt, kind)

    # =========================================================================
    #  Document Generators (PPT, Word, Excel)
    # =========================================================================

    @staticmethod
    async def generate_ppt(
        prompt: str,
        base_url: str,
        provider: str = "auto",
        tier: int = 1,
        attachments: list[dict[str, Any]] | None = None,
        attachment_urls: list[str] | None = None,
    ) -> str:
        """Generate a PPT presentation based on prompt."""
        full_prompt = DocumentService.prepare_prompt_with_attachments(
            prompt, attachments, attachment_urls
        )
        system_prompt = (
            "You are a professional presentation planner. Output ONLY a valid JSON object with a single key 'slides' containing an array of objects. "
            "Do not include any markdown format, backticks, or text before/after the JSON.\n"
            "Format:\n"
            "{\n"
            '  "slides": [\n'
            "    {\n"
            '      "title": "Slide Title",\n'
            '      "bullet_points": ["Point 1", "Point 2", "Point 3"]\n'
            "    }\n"
            "  ]\n"
            "}"
        )

        data = await DocumentService._generate_content(
            full_prompt, system_prompt, provider, tier, kind="ppt"
        )
        slides = DocumentService._normalize_slides(data, prompt)

        prs = Presentation()
        # Set to standard wide 16:9 aspect ratio
        prs.slide_width = 12192000
        prs.slide_height = 6858000

        for slide_data in slides:
            # 1 is Title + Content layout
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # Populate title
            slide.shapes.title.text = slide_data.get("title", "Untitled Slide")

            # Populate body/bullets
            placeholders = slide.placeholders
            if len(placeholders) > 1:
                tf = placeholders[1].text_frame
                tf.text = ""
                for i, pt in enumerate(slide_data.get("bullet_points", [])):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = pt
                    p.level = 0

        attachment_rows = DocumentService._non_image_attachment_rows(attachments)
        if attachment_rows or attachment_urls:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Uploaded Attachments"
            tf = slide.placeholders[1].text_frame
            tf.text = ""
            rows = attachment_rows + [
                {
                    "filename": url,
                    "kind": "url",
                    "mime_type": "link",
                    "size": "",
                    "path": "",
                    "url": url,
                }
                for url in attachment_urls or []
            ]
            for index, row in enumerate(rows):
                p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
                p.text = f"{row['kind'].title()}: {row['filename']}"
                p.level = 0

        for image in DocumentService._image_attachments(attachments):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.25), PptInches(12.3), PptInches(0.5))
            title_frame = title_box.text_frame
            title_frame.text = image["filename"]
            title_frame.paragraphs[0].font.size = Pt(20)
            try:
                slide.shapes.add_picture(
                    image["path"],
                    PptInches(0.7),
                    PptInches(1.0),
                    width=PptInches(8.8),
                    height=None,
                )
            except Exception:
                note = slide.shapes.add_textbox(PptInches(0.7), PptInches(1.2), PptInches(11.0), PptInches(1.0))
                note.text_frame.text = f"Image attachment could not be embedded: {image['filename']}"

        # Save presentation
        os.makedirs("static/generated", exist_ok=True)
        filename = f"ppt_{uuid.uuid4().hex[:12]}.pptx"
        file_path = os.path.join("static/generated", filename)
        prs.save(file_path)

        return DocumentService._save_generated_file(filename, file_path, base_url)

    @staticmethod
    async def generate_word(
        prompt: str,
        base_url: str,
        provider: str = "auto",
        tier: int = 1,
        attachments: list[dict[str, Any]] | None = None,
        attachment_urls: list[str] | None = None,
    ) -> str:
        """Generate a Word document based on prompt."""
        full_prompt = DocumentService.prepare_prompt_with_attachments(
            prompt, attachments, attachment_urls
        )
        system_prompt = (
            "You are a professional document writer. Output ONLY a valid JSON object representing the document outline. "
            "Do not include any markdown format, backticks, or text before/after the JSON. "
            "Format:\n"
            "{\n"
            '  "title": "Document Title",\n'
            '  "sections": [\n'
            "    {\n"
            '      "heading": "Section Title",\n'
            '      "paragraphs": ["First paragraph text.", "Second paragraph text."]\n'
            "    }\n"
            "  ]\n"
            "}"
        )

        data = await DocumentService._generate_content(
            full_prompt, system_prompt, provider, tier, kind="word"
        )
        data = DocumentService._normalize_sections(data, prompt, "word")

        doc = Document()
        doc.add_heading(data.get("title", "Document"), level=0)

        for section in data.get("sections", []):
            doc.add_heading(section.get("heading", ""), level=1)
            for p_text in section.get("paragraphs", []):
                doc.add_paragraph(p_text)

        image_rows = DocumentService._image_attachments(attachments)
        if image_rows:
            doc.add_heading("Uploaded Images", level=1)
            for image in image_rows:
                doc.add_paragraph(image["filename"])
                try:
                    doc.add_picture(image["path"], width=DocxInches(5.8))
                except Exception:
                    doc.add_paragraph(f"Image attachment could not be embedded: {image['filename']}")

        reference_rows = DocumentService._non_image_attachment_rows(attachments)
        if reference_rows or attachment_urls:
            doc.add_heading("Uploaded Attachments", level=1)
            for row in reference_rows:
                doc.add_paragraph(
                    f"{row['kind'].title()}: {row['filename']} ({row['mime_type']})",
                    style="List Bullet",
                )
            for url in attachment_urls or []:
                doc.add_paragraph(f"URL: {url}", style="List Bullet")

        # Save document
        os.makedirs("static/generated", exist_ok=True)
        filename = f"doc_{uuid.uuid4().hex[:12]}.docx"
        file_path = os.path.join("static/generated", filename)
        doc.save(file_path)

        return DocumentService._save_generated_file(filename, file_path, base_url)

    @staticmethod
    async def generate_excel(
        prompt: str,
        base_url: str,
        provider: str = "auto",
        tier: int = 1,
        attachments: list[dict[str, Any]] | None = None,
        attachment_urls: list[str] | None = None,
    ) -> str:
        """Generate an Excel spreadsheet based on prompt."""
        full_prompt = DocumentService.prepare_prompt_with_attachments(
            prompt, attachments, attachment_urls
        )
        system_prompt = (
            "You are an expert database analyst. Output ONLY a valid JSON object representing a data spreadsheet. "
            "Do not include any markdown format, backticks, or text before/after the JSON. "
            "Format:\n"
            "{\n"
            '  "sheet_title": "Sheet Name",\n'
            '  "headers": ["Column 1", "Column 2", "Column 3"],\n'
            '  "rows": [\n'
            '    ["Value 1a", "Value 1b", "Value 1c"],\n'
            '    ["Value 2a", "Value 2b", "Value 2c"]\n'
            "  ]\n"
            "}"
        )

        data = await DocumentService._generate_content(
            full_prompt, system_prompt, provider, tier, kind="excel"
        )
        data = DocumentService._normalize_sheet(data, prompt)

        wb = Workbook()
        ws = wb.active
        ws.title = data.get("sheet_title", "DataSheet")[:30]  # Excel limit is 31 chars

        # Append headers
        ws.append(data.get("headers", []))

        # Append rows
        for row in data.get("rows", []):
            ws.append(row)

        # Style header row (optional but nice)
        for col_num in range(1, len(data.get("headers", [])) + 1):
            cell = ws.cell(row=1, column=col_num)
            cell.font = cell.font.copy(bold=True)

        attachment_rows = DocumentService._attachment_rows(attachments)
        if attachment_rows or attachment_urls:
            attach_ws = wb.create_sheet("Attachments")
            attach_ws.append(["Type", "Filename or URL", "MIME Type", "Size Bytes"])
            for cell in attach_ws[1]:
                cell.font = cell.font.copy(bold=True)

            for row in attachment_rows:
                attach_ws.append([row["kind"], row["filename"], row["mime_type"], row["size"]])

            for url in attachment_urls or []:
                attach_ws.append(["url", url, "link", ""])

            image_row = len(attachment_rows) + len(attachment_urls or []) + 4
            for image in DocumentService._image_attachments(attachments):
                try:
                    excel_image = ExcelImage(image["path"])
                    excel_image.width = min(excel_image.width, 480)
                    excel_image.height = min(excel_image.height, 320)
                    attach_ws.add_image(excel_image, f"A{image_row}")
                    attach_ws.cell(row=image_row, column=4, value=image["filename"])
                    image_row += 18
                except Exception:
                    attach_ws.append(["image", image["filename"], "embed_failed", image["size"]])

        # Save spreadsheet
        os.makedirs("static/generated", exist_ok=True)
        filename = f"excel_{uuid.uuid4().hex[:12]}.xlsx"
        file_path = os.path.join("static/generated", filename)
        wb.save(file_path)

        return DocumentService._save_generated_file(filename, file_path, base_url)

    # =========================================================================
    #  PDF Document Generator (fpdf2 – Professional Indigo Theme)
    # =========================================================================

    @staticmethod
    async def generate_pdf(
        prompt: str,
        base_url: str,
        provider: str = "auto",
        tier: int = 1,
        attachments: list[dict[str, Any]] | None = None,
        attachment_urls: list[str] | None = None,
    ) -> str:
        """Generate a professionally styled PDF document based on prompt using AI-structured content."""
        from fpdf import FPDF

        full_prompt = DocumentService.prepare_prompt_with_attachments(
            prompt, attachments, attachment_urls
        )

        system_prompt = (
            "You are a professional document writer. Output ONLY a valid JSON object representing a document outline. "
            "Do not include any markdown format, backticks, or text before/after the JSON.\n"
            "Format:\n"
            "{\n"
            '  "title": "Document Title",\n'
            '  "subtitle": "A brief one-line description",\n'
            '  "sections": [\n'
            "    {\n"
            '      "heading": "Section Title",\n'
            '      "paragraphs": ["First paragraph text.", "Second paragraph text."]\n'
            "    }\n"
            "  ]\n"
            "}"
        )

        data = await DocumentService._generate_content(
            full_prompt, system_prompt, provider, tier, kind="pdf"
        )
        data = DocumentService._normalize_sections(data, prompt, "pdf")

        # ── Theme Constants ──────────────────────────────────────────
        INDIGO = (79, 70, 229)  # Primary accent
        CHARCOAL = (31, 41, 55)  # Body text
        LIGHT_GRAY = (243, 244, 246)  # Subtle background fills
        WHITE = (255, 255, 255)
        DARK_GRAY = (107, 114, 128)  # Footer / meta text

        class VedaPDF(FPDF):
            """Custom FPDF subclass that draws themed headers and page-numbered footers."""

            def __init__(self, doc_title: str):
                super().__init__()
                self.doc_title = doc_title

            def header(self):
                # Indigo header band
                self.set_fill_color(*INDIGO)
                self.rect(0, 0, self.w, 18, "F")
                self.set_font("Helvetica", "B", 11)
                self.set_text_color(*WHITE)
                self.set_xy(self.l_margin, 4)
                self.cell(self.epw, 10, DocumentService._pdf_safe_text(self.doc_title[:80]), align="L")

                # Timestamp on the right
                self.set_font("Helvetica", "", 8)
                self.set_xy(self.w - self.r_margin - 50, 5)
                self.cell(50, 8, utcnow().strftime("%Y-%m-%d %H:%M UTC"), align="R")
                self.set_y(22)

            def footer(self):
                self.set_y(-15)
                self.set_font("Helvetica", "I", 8)
                self.set_text_color(*DARK_GRAY)
                self.cell(0, 10, f"Page {self.page_no()}/{{nb}}", align="C")

        doc_title = data.get("title", "Document")
        doc_subtitle = data.get("subtitle", "")
        sections = data.get("sections", [])
        if not isinstance(sections, list):
            sections = [data]

        pdf = VedaPDF(doc_title)
        pdf.alias_nb_pages()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()

        # ── Cover Title Block ────────────────────────────────────────
        pdf.set_y(30)
        pdf.set_font("Helvetica", "B", 26)
        pdf.set_text_color(*INDIGO)
        pdf.multi_cell(pdf.epw, 12, DocumentService._pdf_safe_text(doc_title), align="C")
        pdf.ln(4)

        if doc_subtitle:
            pdf.set_font("Helvetica", "", 12)
            pdf.set_text_color(*DARK_GRAY)
            pdf.multi_cell(pdf.epw, 8, DocumentService._pdf_safe_text(doc_subtitle), align="C")
            pdf.ln(2)

        # Accent divider line
        pdf.ln(6)
        y_line = pdf.get_y()
        pdf.set_draw_color(*INDIGO)
        pdf.set_line_width(0.8)
        pdf.line(pdf.l_margin + 20, y_line, pdf.w - pdf.r_margin - 20, y_line)
        pdf.ln(10)

        # ── Body Sections ────────────────────────────────────────────
        for idx, section in enumerate(sections):
            heading = section.get("heading", f"Section {idx + 1}")
            paragraphs = section.get("paragraphs", [])

            # Section heading with indigo left-bar accent
            y_before = pdf.get_y()
            if y_before > 250:  # Slightly more conservative page break
                pdf.add_page()
                y_before = pdf.get_y()

            pdf.set_fill_color(*INDIGO)
            pdf.rect(pdf.l_margin, y_before, 2.5, 9, "F")
            pdf.set_xy(pdf.l_margin + 6, y_before)
            pdf.set_font("Helvetica", "B", 14)
            pdf.set_text_color(*CHARCOAL)
            pdf.cell(pdf.epw - 6, 9, DocumentService._pdf_safe_text(heading))
            pdf.ln(12)

            # Paragraphs
            pdf.set_font("Helvetica", "", 11)
            pdf.set_text_color(*CHARCOAL)
            for para in paragraphs:
                if not isinstance(para, str):
                    para = str(para)
                pdf.multi_cell(pdf.epw, 6.5, DocumentService._pdf_safe_text(para))
                pdf.ln(3)

            pdf.ln(4)

        reference_rows = DocumentService._non_image_attachment_rows(attachments)
        if reference_rows or attachment_urls:
            pdf.add_page()
            pdf.set_font("Helvetica", "B", 16)
            pdf.set_text_color(*CHARCOAL)
            pdf.cell(0, 10, "Uploaded Attachments")
            pdf.ln(12)
            pdf.set_font("Helvetica", "", 10)
            for row in reference_rows:
                line = f"{row['kind'].title()}: {row['filename']} ({row['mime_type']})"
                pdf.multi_cell(pdf.epw, 6, DocumentService._pdf_safe_text(line))
            for url in attachment_urls or []:
                pdf.multi_cell(pdf.epw, 6, DocumentService._pdf_safe_text(f"URL: {url}"))

        for image in DocumentService._image_attachments(attachments):
            pdf.add_page()
            pdf.set_font("Helvetica", "B", 14)
            pdf.set_text_color(*CHARCOAL)
            pdf.multi_cell(pdf.epw, 8, DocumentService._pdf_safe_text(image["filename"]))
            pdf.ln(4)
            try:
                pdf.image(image["path"], x=pdf.l_margin, w=min(pdf.epw, 170))
            except Exception:
                pdf.set_font("Helvetica", "", 10)
                pdf.multi_cell(
                    pdf.epw,
                    6,
                    DocumentService._pdf_safe_text(
                        f"Image attachment could not be embedded: {image['filename']}"
                    ),
                )

        # ── Save PDF ─────────────────────────────────────────────────
        os.makedirs("static/generated", exist_ok=True)
        filename = f"pdf_{uuid.uuid4().hex[:12]}.pdf"
        file_path = os.path.join("static/generated", filename)
        pdf.output(file_path)

        return DocumentService._save_generated_file(filename, file_path, base_url)
