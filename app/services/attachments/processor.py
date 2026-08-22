from pathlib import Path
from typing import Optional

from .config import ATTACHMENT_CONFIG
from .models import AttachmentMetadata, ParsedAttachment


class AttachmentProcessor:
    @staticmethod
    def _truncate(text: str) -> str:
        return text[: ATTACHMENT_CONFIG.MAX_DOCUMENT_CHARACTERS]

    @staticmethod
    def read_file_bytes(path: str) -> bytes:
        return Path(path).read_bytes()

    @staticmethod
    def extract_document_text(path: str, mime_type: str) -> Optional[str]:
        file_path = Path(path)
        extension = file_path.suffix.lower()

        if mime_type == "application/pdf":
            try:
                from pypdf import PdfReader

                reader = PdfReader(str(file_path))
                pages = []
                for page in reader.pages[:10]:
                    pages.append(page.extract_text() or "")
                return AttachmentProcessor._truncate("\n".join(pages).strip())
            except Exception:
                return "[PDF attachment provided; text extraction is not available in this environment.]"

        if mime_type in {"text/plain", "text/csv", "application/json"} or extension in {
            ".txt",
            ".csv",
            ".json",
        }:
            raw = file_path.read_bytes()
            text = raw.decode("utf-8", errors="replace")
            return AttachmentProcessor._truncate(text)

        if (
            mime_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or extension == ".docx"
        ):
            try:
                from docx import Document

                doc = Document(str(file_path))
                parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
                for table in doc.tables[:5]:
                    for row in table.rows[:20]:
                        parts.append(" | ".join(cell.text.strip() for cell in row.cells))
                return AttachmentProcessor._truncate("\n".join(parts))
            except Exception:
                return "[Word attachment provided; text extraction failed.]"

        if (
            mime_type
            == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            or extension == ".xlsx"
        ):
            try:
                from openpyxl import load_workbook

                wb = load_workbook(str(file_path), read_only=True, data_only=True)
                parts = []
                for sheet in wb.worksheets[:5]:
                    parts.append(f"Sheet: {sheet.title}")
                    for row in sheet.iter_rows(max_row=30, max_col=12, values_only=True):
                        values = ["" if value is None else str(value) for value in row]
                        if any(values):
                            parts.append(" | ".join(values))
                wb.close()
                return AttachmentProcessor._truncate("\n".join(parts))
            except Exception:
                return "[Excel attachment provided; text extraction failed.]"

        if (
            mime_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or extension == ".pptx"
        ):
            try:
                from pptx import Presentation

                presentation = Presentation(str(file_path))
                parts = []
                for index, slide in enumerate(presentation.slides[:20], start=1):
                    slide_text = []
                    for shape in slide.shapes:
                        text = getattr(shape, "text", "")
                        if text and text.strip():
                            slide_text.append(text.strip())
                    if slide_text:
                        parts.append(f"Slide {index}: " + " | ".join(slide_text))
                return AttachmentProcessor._truncate("\n".join(parts))
            except Exception:
                return "[PowerPoint attachment provided; text extraction failed.]"

        return None

    @staticmethod
    def parse_attachment(metadata: AttachmentMetadata) -> ParsedAttachment:
        content = Path(metadata.temp_path).read_bytes()
        text_preview = None
        if metadata.is_document:
            text_preview = AttachmentProcessor.extract_document_text(metadata.temp_path, metadata.mime_type)
        return ParsedAttachment(
            attachment=metadata,
            content=content,
            text_preview=text_preview,
        )
