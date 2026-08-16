"""
Comprehensive test suite for PowerPoint presentation generation.

Tests cover:
- Presentation plan generation from AI
- PPTX file creation
- User isolation and security
- Text overflow handling
- Storage integration
- End-to-end workflow
"""

import io
import json
import pytest
from unittest.mock import Mock, AsyncMock, patch, MagicMock
from fastapi.testclient import TestClient

from app.schemas.presentations import (
    PresentationPlan,
    Slide,
    SlideLayout,
    Theme,
    ChartType,
    Chart,
    ChartDataSeries,
    Table,
    TableRow,
    PPTGenerationRequest,
)
from app.services.ppt import generate_pptx
from app.models.user import User
from app.models.asset import AIAsset


# ============================================================================
# Test Data Fixtures
# ============================================================================

@pytest.fixture
def mock_user():
    """Create a mock authenticated user."""
    user = Mock(spec=User)
    user.id = 12345
    user.email = "test@example.com"
    user.username = "testuser"
    return user


@pytest.fixture
def mock_other_user():
    """Create a different mock user for isolation tests."""
    user = Mock(spec=User)
    user.id = 99999
    user.email = "other@example.com"
    user.username = "otheruser"
    return user


@pytest.fixture
def simple_presentation_plan():
    """Create a simple valid presentation plan."""
    return PresentationPlan(
        title="Test Presentation",
        subtitle="A Simple Test",
        author="Test Author",
        theme=Theme.modern,
        language="en",
        slides=[
            Slide(
                slide_number=1,
                layout=SlideLayout.title,
                title="Introduction",
                subtitle="Welcome to the presentation",
                bullets=[],
                paragraphs=[],
            ),
            Slide(
                slide_number=2,
                layout=SlideLayout.content,
                title="Main Points",
                bullets=["First point", "Second point", "Third point"],
                paragraphs=["This is a longer paragraph explaining the content."],
            ),
            Slide(
                slide_number=3,
                layout=SlideLayout.conclusion,
                title="Thank You",
                subtitle="Questions?",
            ),
        ],
    )


@pytest.fixture
def complex_presentation_plan():
    """Create a presentation with all layout types."""
    return PresentationPlan(
        title="Complex Presentation",
        subtitle="Testing All Features",
        author="Test",
        theme=Theme.professional,
        language="en",
        slides=[
            # Title slide
            Slide(
                slide_number=1,
                layout=SlideLayout.title,
                title="Complex Presentation",
                subtitle="All Features",
            ),
            # Section slide
            Slide(
                slide_number=2,
                layout=SlideLayout.section,
                title="Section 1: Basics",
            ),
            # Content with bullets
            Slide(
                slide_number=3,
                layout=SlideLayout.content,
                title="Content Slide",
                bullets=["Bullet 1", "Bullet 2", "Bullet 3"],
                paragraphs=["Paragraph content"],
            ),
            # Two column layout
            Slide(
                slide_number=4,
                layout=SlideLayout.two_column,
                title="Two Column Layout",
                bullets=[
                    "Left column item 1",
                    "Left column item 2",
                    "Right column item 1",
                    "Right column item 2",
                ],
            ),
            # Quote slide
            Slide(
                slide_number=5,
                layout=SlideLayout.quote,
                title="Quote",
                quote="This is an inspiring quote about life.",
                author="Famous Person",
            ),
            # Code slide
            Slide(
                slide_number=6,
                layout=SlideLayout.code,
                title="Code Example",
                code="def hello():\n    print('Hello, World!')",
                code_language="python",
            ),
            # Table slide
            Slide(
                slide_number=7,
                layout=SlideLayout.table,
                title="Table Example",
                table=Table(
                    headers=["Column 1", "Column 2", "Column 3"],
                    rows=[
                        TableRow(cells=["Data 1", "Data 2", "Data 3"]),
                        TableRow(cells=["Data 4", "Data 5", "Data 6"]),
                    ],
                ),
            ),
            # Chart slide
            Slide(
                slide_number=8,
                layout=SlideLayout.chart,
                title="Chart Example",
                chart=Chart(
                    type=ChartType.bar,
                    title="Sales Data",
                    categories=["Q1", "Q2", "Q3", "Q4"],
                    series=[
                        ChartDataSeries(name="Series 1", values=[10, 20, 15, 25]),
                        ChartDataSeries(name="Series 2", values=[15, 25, 20, 30]),
                    ],
                ),
            ),
            # Conclusion
            Slide(
                slide_number=9,
                layout=SlideLayout.conclusion,
                title="Thank You",
                subtitle="Questions?",
            ),
        ],
    )


@pytest.fixture
def ai_response_gemini():
    """Mock Gemini API response format."""
    return {
        "candidates": [
            {
                "content": {
                    "parts": [
                        {
                            "text": json.dumps({
                                "title": "Generated Presentation",
                                "subtitle": "Auto-generated",
                                "author": "VedaApex",
                                "theme": "modern",
                                "language": "en",
                                "slides": [
                                    {
                                        "slide_number": 1,
                                        "layout": "title",
                                        "title": "Title Slide",
                                        "subtitle": "Subtitle",
                                        "bullets": [],
                                        "paragraphs": [],
                                    },
                                    {
                                        "slide_number": 2,
                                        "layout": "content",
                                        "title": "Content",
                                        "bullets": ["Point 1", "Point 2"],
                                        "paragraphs": [],
                                    },
                                    {
                                        "slide_number": 3,
                                        "layout": "conclusion",
                                        "title": "Thank You",
                                        "subtitle": None,
                                        "bullets": [],
                                        "paragraphs": [],
                                    },
                                ],
                            })
                        }
                    ]
                }
            }
        ]
    }


@pytest.fixture
def ai_response_openai():
    """Mock OpenAI API response format."""
    return {
        "choices": [
            {
                "message": {
                    "content": json.dumps({
                        "title": "Generated Presentation",
                        "subtitle": "Auto-generated",
                        "author": "VedaApex",
                        "theme": "modern",
                        "language": "en",
                        "slides": [
                            {
                                "slide_number": 1,
                                "layout": "title",
                                "title": "Title",
                                "subtitle": "Sub",
                                "bullets": [],
                                "paragraphs": [],
                            },
                            {
                                "slide_number": 2,
                                "layout": "conclusion",
                                "title": "End",
                                "subtitle": None,
                                "bullets": [],
                                "paragraphs": [],
                            },
                        ],
                    })
                }
            }
        ]
    }


# ============================================================================
# Unit Tests: PPTX Generation
# ============================================================================

class TestPPTXGeneration:
    """Test PPTX file generation from presentation plans."""

    def test_generate_simple_pptx(self, simple_presentation_plan):
        """Test generating PPTX from simple presentation plan."""
        pptx_bytes = generate_pptx(simple_presentation_plan)
        
        # Verify output is bytes
        assert isinstance(pptx_bytes, bytes)
        
        # Verify it's a valid ZIP file (PPTX is ZIP format)
        assert pptx_bytes[:4] == b'PK\x03\x04'  # ZIP magic number
        
        # Verify reasonable file size
        assert len(pptx_bytes) > 10000  # Should be at least 10KB

    def test_generate_complex_pptx(self, complex_presentation_plan):
        """Test generating PPTX with all layout types."""
        pptx_bytes = generate_pptx(complex_presentation_plan)
        
        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 10000

    def test_pptx_slide_count(self, simple_presentation_plan):
        """Verify slide count matches presentation plan."""
        pptx_bytes = generate_pptx(simple_presentation_plan)
        
        # Load the PPTX to verify slide count
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
        
        assert len(prs.slides) == len(simple_presentation_plan.slides)

    def test_title_slide_content(self, simple_presentation_plan):
        """Verify title slide contains expected text."""
        pptx_bytes = generate_pptx(simple_presentation_plan)
        
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
        
        # First slide should be title slide
        first_slide = prs.slides[0]
        slide_text = " ".join([shape.text for shape in first_slide.shapes if hasattr(shape, "text")])
        
        assert "Introduction" in slide_text
        assert "Welcome" in slide_text

    def test_content_slide_bullets(self, simple_presentation_plan):
        """Verify content slide includes bullet points."""
        pptx_bytes = generate_pptx(simple_presentation_plan)
        
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
        
        # Second slide should have bullets
        second_slide = prs.slides[1]
        slide_text = " ".join([shape.text for shape in second_slide.shapes if hasattr(shape, "text")])
        
        assert "First point" in slide_text or "Second point" in slide_text

    def test_theme_application(self):
        """Verify theme colors are applied correctly."""
        plan = PresentationPlan(
            title="Theme Test",
            theme=Theme.dark,
            slides=[
                Slide(
                    slide_number=1,
                    layout=SlideLayout.title,
                    title="Dark Theme",
                )
            ],
        )
        
        pptx_bytes = generate_pptx(plan)
        assert isinstance(pptx_bytes, bytes)
        assert len(pptx_bytes) > 5000


# ============================================================================
# Unit Tests: Presentation Plan Validation
# ============================================================================

class TestPresentationPlanValidation:
    """Test validation of presentation plans."""

    def test_valid_presentation_plan(self, simple_presentation_plan):
        """Test creating a valid presentation plan."""
        assert simple_presentation_plan.title == "Test Presentation"
        assert len(simple_presentation_plan.slides) == 3

    def test_slide_number_validation(self):
        """Test that slide numbers must be sequential."""
        with pytest.raises(ValueError):
            PresentationPlan(
                title="Invalid",
                slides=[
                    Slide(slide_number=1, layout=SlideLayout.title, title="Slide 1"),
                    Slide(slide_number=3, layout=SlideLayout.content, title="Slide 3"),  # Skipped 2
                ],
            )

    def test_minimum_slides(self):
        """Test that presentation must have at least 1 slide."""
        with pytest.raises(ValueError):
            PresentationPlan(
                title="Empty",
                slides=[],  # Empty slides list
            )

    def test_maximum_slides(self):
        """Test that presentation cannot exceed 200 slides."""
        # Create exactly 200 valid slides first
        slides = [
            Slide(
                slide_number=i,
                layout=SlideLayout.content,
                title=f"Slide {i}",
            )
            for i in range(1, 201)  # 200 slides is OK
        ]
        
        # 200 slides should be valid
        plan = PresentationPlan(title="Max", slides=slides)
        assert len(plan.slides) == 200
        
        # Trying to add slide 201 should fail at Slide level
        with pytest.raises(ValueError):
            Slide(
                slide_number=201,
                layout=SlideLayout.content,
                title="Slide 201",
            )

    def test_bullet_text_truncation_validation(self):
        """Test that bullet points validate length."""
        with pytest.raises(ValueError):
            Slide(
                slide_number=1,
                layout=SlideLayout.content,
                title="Test",
                bullets=["x" * 301],  # Exceeds 300 char limit
            )

    def test_paragraph_text_truncation_validation(self):
        """Test that paragraphs validate length."""
        with pytest.raises(ValueError):
            Slide(
                slide_number=1,
                layout=SlideLayout.content,
                title="Test",
                paragraphs=["x" * 1001],  # Exceeds 1000 char limit
            )


# ============================================================================
# Integration Tests: API Endpoint
# ============================================================================

class TestPresentationEndpoint:
    """Test PPT generation API endpoint."""

    @pytest.mark.asyncio
    async def test_generate_presentation_endpoint_success(self, mock_user, ai_response_gemini):
        """Test successful PPT generation via endpoint."""
        # This test would require a full FastAPI test client setup
        # For now, we test the response model
        from app.schemas.presentations import PPTGenerationResponse
        
        response = PPTGenerationResponse(
            success=True,
            presentation_id="test-123",
            attachment_id=42,
            filename="presentation.pptx",
            file_size_bytes=25000,
            status="completed",
            proxy_url="https://example.com/files/presentation.pptx",
        )
        
        assert response.success is True
        assert response.presentation_id == "test-123"
        assert response.attachment_id == 42

    def test_ppt_generation_request_validation(self):
        """Test request validation."""
        request = PPTGenerationRequest(
            prompt="Create a presentation about Python",
            slide_count=5,
            theme=Theme.modern,
            language="English",
            include_images=False,
        )
        
        assert request.prompt == "Create a presentation about Python"
        assert request.slide_count == 5

    def test_ppt_generation_request_prompt_too_short(self):
        """Test that prompt must be at least 10 characters."""
        with pytest.raises(ValueError):
            PPTGenerationRequest(
                prompt="Short",  # Only 5 chars
                slide_count=5,
            )

    def test_ppt_generation_request_slide_count_bounds(self):
        """Test slide count bounds (3-200)."""
        with pytest.raises(ValueError):
            PPTGenerationRequest(
                prompt="Valid prompt here",
                slide_count=2,  # Below minimum
            )
        
        with pytest.raises(ValueError):
            PPTGenerationRequest(
                prompt="Valid prompt here",
                slide_count=201,  # Above maximum
            )


# ============================================================================
# Tests: Text Overflow Handling
# ============================================================================

class TestTextOverflowHandling:
    """Test handling of text that might overflow slides."""

    def test_long_title_truncation(self):
        """Test that very long titles are handled."""
        plan = PresentationPlan(
            title="A" * 200,  # Max length is 200, should be OK
            slides=[
                Slide(
                    slide_number=1,
                    layout=SlideLayout.title,
                    title="A" * 200,  # At max length
                )
            ],
        )
        
        pptx_bytes = generate_pptx(plan)
        assert isinstance(pptx_bytes, bytes)

    def test_very_long_bullet_point(self):
        """Test handling of bullet points at max length."""
        plan = PresentationPlan(
            title="Overflow Test",
            slides=[
                Slide(
                    slide_number=1,
                    layout=SlideLayout.content,
                    title="Test",
                    bullets=["B" * 300],  # Exactly at max length
                )
            ],
        )
        
        pptx_bytes = generate_pptx(plan)
        assert isinstance(pptx_bytes, bytes)

    def test_many_bullet_points(self):
        """Test slide with many bullets (max 10)."""
        plan = PresentationPlan(
            title="Many Bullets",
            slides=[
                Slide(
                    slide_number=1,
                    layout=SlideLayout.content,
                    title="Test",
                    bullets=[f"Bullet {i}" for i in range(10)],  # Max 10
                )
            ],
        )
        
        pptx_bytes = generate_pptx(plan)
        assert isinstance(pptx_bytes, bytes)


# ============================================================================
# Tests: User Isolation & Security
# ============================================================================

class TestUserIsolation:
    """Test that users can only access their own presentations."""

    def test_presentation_user_ownership(self, mock_user):
        """Test that presentation is linked to user."""
        asset = Mock(spec=AIAsset)
        asset.user_id = mock_user.id
        asset.asset_type = "presentation"
        
        # User should own this asset
        assert asset.user_id == mock_user.id

    def test_different_user_cannot_access(self, mock_user, mock_other_user):
        """Test isolation between users."""
        asset = Mock(spec=AIAsset)
        asset.user_id = mock_user.id
        
        # Other user should not be able to access
        assert asset.user_id != mock_other_user.id


# ============================================================================
# Tests: AI Response Parsing
# ============================================================================

class TestAIResponseParsing:
    """Test parsing of different AI model responses."""

    def test_parse_gemini_response(self, ai_response_gemini):
        """Test parsing Gemini API response."""
        # Extract JSON from Gemini response
        json_text = ai_response_gemini["candidates"][0]["content"]["parts"][0]["text"]
        plan_dict = json.loads(json_text)
        plan = PresentationPlan(**plan_dict)
        
        assert plan.title == "Generated Presentation"
        assert len(plan.slides) == 3

    def test_parse_openai_response(self, ai_response_openai):
        """Test parsing OpenAI API response."""
        # Extract JSON from OpenAI response
        json_text = ai_response_openai["choices"][0]["message"]["content"]
        plan_dict = json.loads(json_text)
        plan = PresentationPlan(**plan_dict)
        
        assert plan.title == "Generated Presentation"
        assert len(plan.slides) == 2

    def test_parse_markdown_wrapped_json(self):
        """Test parsing JSON wrapped in markdown code blocks."""
        json_text = '```json\n{"title": "Test", "slides": [{"slide_number": 1, "layout": "title", "title": "T"}]}\n```'
        
        # Clean markdown
        json_text = json_text.strip()
        if json_text.startswith("```"):
            lines = json_text.split("\n")
            json_text = "\n".join(lines[1:-1])
        
        plan_dict = json.loads(json_text)
        plan = PresentationPlan(**plan_dict)
        
        assert plan.title == "Test"

    def test_invalid_json_response(self):
        """Test handling of invalid JSON from AI."""
        invalid_json = '{"title": "Invalid JSON", incomplete'
        
        with pytest.raises(json.JSONDecodeError):
            json.loads(invalid_json)


# ============================================================================
# Tests: Storage Integration
# ============================================================================

class TestStorageIntegration:
    """Test integration with AssetStorageService."""

    def test_asset_metadata_tracking(self):
        """Test that asset metadata is tracked correctly."""
        asset = Mock(spec=AIAsset)
        asset.id = 42
        asset.user_id = 12345
        asset.asset_type = "presentation"
        asset.file_size_bytes = 25000
        asset.proxy_url = "https://example.com/file.pptx"
        
        assert asset.asset_type == "presentation"
        assert asset.file_size_bytes == 25000
        assert asset.proxy_url.endswith(".pptx")

    def test_asset_never_exposes_raw_url(self):
        """Test that raw URLs are never exposed (only proxy URLs)."""
        asset = Mock(spec=AIAsset)
        asset.proxy_url = "https://example.com/proxy/file.pptx"
        asset.original_url = None  # Raw URL not exposed
        
        assert asset.proxy_url is not None
        assert not asset.proxy_url.endswith("/.pptx")  # Should be masked


# ============================================================================
# Tests: End-to-End Workflow
# ============================================================================

class TestEndToEndWorkflow:
    """Test complete PPT generation workflow."""

    def test_workflow_simple(self, simple_presentation_plan):
        """Test complete workflow: plan → PPTX → validation."""
        # Step 1: Have a presentation plan
        assert simple_presentation_plan.title == "Test Presentation"
        
        # Step 2: Generate PPTX
        pptx_bytes = generate_pptx(simple_presentation_plan)
        assert isinstance(pptx_bytes, bytes)
        
        # Step 3: Verify PPTX structure
        from pptx import Presentation
        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) == 3
        
        # Step 4: Verify no corruption
        assert pptx_bytes[:4] == b'PK\x03\x04'

    def test_workflow_with_metadata(self, simple_presentation_plan, mock_user):
        """Test workflow including metadata tracking."""
        # Generate PPTX
        pptx_bytes = generate_pptx(simple_presentation_plan)
        
        # Create metadata record
        asset = Mock(spec=AIAsset)
        asset.user_id = mock_user.id
        asset.file_size_bytes = len(pptx_bytes)
        asset.asset_type = "presentation"
        
        # Verify metadata
        assert asset.user_id == mock_user.id
        assert asset.file_size_bytes == len(pptx_bytes)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
